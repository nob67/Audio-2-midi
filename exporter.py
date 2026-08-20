import os
from PIL import Image
from reportlab.lib.pagesizes import letter, landscape
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches, Pt


def export_to_pdf(layout_data, output_path, page_size=landscape(letter), margin=36, spacing=18):
    """
    Exports layout data to a PDF file using ReportLab.

    Args:
        layout_data: dict produced by calculate_imposition
        output_path: file path for output PDF
        page_size: tuple (width, height) in points
        margin: page margin in points (1 point = 1/72 inch)
        spacing: spacing between grid items in points
    """
    page_width, page_height = page_size
    c = canvas.Canvas(output_path, pagesize=page_size)

    rows = layout_data['rows']
    cols = layout_data['cols']

    avail_width = page_width - (2 * margin) - ((cols - 1) * spacing)
    avail_height = page_height - (2 * margin) - ((rows - 1) * spacing)

    cell_width = avail_width / cols if cols > 0 else avail_width
    cell_height = avail_height / rows if rows > 0 else avail_height

    for sheet in layout_data['sheets']:
        for side_key in ['front', 'back']:
            side_data = sheet[side_key]
            grid = side_data['grid']
            rotate = side_data.get('rotate_180', False)

            for r in range(rows):
                for col_idx in range(cols):
                    img_path = grid[r][col_idx]
                    if not img_path or not os.path.exists(str(img_path)):
                        continue

                    # Compute bounding box for cell
                    # PDF coordinates start from bottom-left
                    left = margin + col_idx * (cell_width + spacing)
                    top_y = page_height - margin - r * (cell_height + spacing)
                    bottom = top_y - cell_height

                    c.saveState()
                    if rotate:
                        # Rotate 180 degrees around center of cell
                        cx = left + cell_width / 2.0
                        cy = bottom + cell_height / 2.0
                        c.translate(cx, cy)
                        c.rotate(180)
                        c.translate(-cx, -cy)

                    try:
                        with Image.open(img_path) as img:
                            img_w, img_h = img.size
                            scale = min(cell_width / img_w, cell_height / img_h)
                            draw_w = img_w * scale
                            draw_h = img_h * scale

                            # Center in cell
                            draw_x = left + (cell_width - draw_w) / 2.0
                            draw_y = bottom + (cell_height - draw_h) / 2.0

                            c.drawImage(str(img_path), draw_x, draw_y, width=draw_w, height=draw_h, preserveAspectRatio=True)
                    except Exception as e:
                        print(f"Error drawing image {img_path}: {e}")

                    c.restoreState()

            c.showPage()

    c.save()
    return output_path


def export_to_pptx(layout_data, output_path, slide_width_in=11.0, slide_height_in=8.5, margin_in=0.5, spacing_in=0.25):
    """
    Exports layout data to a PPTX file using python-pptx.

    Args:
        layout_data: dict produced by calculate_imposition
        output_path: file path for output PPTX
        slide_width_in: width of slide in inches
        slide_height_in: height of slide in inches
        margin_in: margin in inches
        spacing_in: spacing between grid items in inches
    """
    prs = Presentation()
    prs.slide_width = Inches(slide_width_in)
    prs.slide_height = Inches(slide_height_in)
    blank_slide_layout = prs.slide_layouts[6]

    rows = layout_data['rows']
    cols = layout_data['cols']

    avail_w = slide_width_in - (2 * margin_in) - ((cols - 1) * spacing_in)
    avail_h = slide_height_in - (2 * margin_in) - ((rows - 1) * spacing_in)

    cell_w = avail_w / cols if cols > 0 else avail_w
    cell_h = avail_h / rows if rows > 0 else avail_h

    for sheet in layout_data['sheets']:
        for side_key in ['front', 'back']:
            side_data = sheet[side_key]
            grid = side_data['grid']
            rotate = side_data.get('rotate_180', False)

            slide = prs.slides.add_slide(blank_slide_layout)

            for r in range(rows):
                for col_idx in range(cols):
                    img_path = grid[r][col_idx]
                    if not img_path or not os.path.exists(str(img_path)):
                        continue

                    left = Inches(margin_in + col_idx * (cell_w + spacing_in))
                    top = Inches(margin_in + r * (cell_h + spacing_in))

                    try:
                        with Image.open(img_path) as img:
                            img_w, img_h = img.size
                            scale = min(cell_w / img_w, cell_h / img_h)
                            draw_w = Inches(img_w * scale)
                            draw_h = Inches(img_h * scale)

                            # Center in cell
                            offset_x = Inches((cell_w - (img_w * scale)) / 2.0)
                            offset_y = Inches((cell_h - (img_h * scale)) / 2.0)

                            shape = slide.shapes.add_picture(
                                str(img_path),
                                left + offset_x,
                                top + offset_y,
                                width=draw_w,
                                height=draw_h
                            )

                            if rotate:
                                shape.rotation = 180
                    except Exception as e:
                        print(f"Error adding picture {img_path} to PPTX: {e}")

    prs.save(output_path)
    return output_path
